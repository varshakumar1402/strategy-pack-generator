from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def replace_text_placeholders(slide, data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for key, value in data.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))

def insert_image(slide, placeholder, image_path, width=Inches(4)):
    if not os.path.exists(image_path):
        return

    for shape in slide.shapes:
        if shape.has_text_frame:
            if placeholder in shape.text:
                left = shape.left
                top = shape.top
                slide.shapes._spTree.remove(shape._element)  # remove placeholder box
                slide.shapes.add_picture(image_path, left, top, width=width)
                break

def create_pptx(data, template_path="template.pptx"):
    prs = Presentation(template_path)

    # Slide-by-slide processing
    for slide in prs.slides:
        replace_text_placeholders(slide, data)

        # Insert known image placeholders
        insert_image(slide, "{{flowchart image}}", data.get("flowchart", ""))
        insert_image(slide, "{{streetview_image}}", data.get("streetview", ""))
        insert_image(slide, "{{Conceptual Floor Plan}}", data.get("conceptual_plan", ""))
        insert_image(slide, "{{Interior image 1}}", data.get("interior_1", ""))
        insert_image(slide, "{{Interior image 2}}", data.get("interior_2", ""))
        insert_image(slide, "{{property - sample 1 - screenshot}}", data.get("sample1_streetview", ""))
        insert_image(slide, "{{property - sample 2 - screenshot}}", data.get("sample2_streetview", ""))

    # Save final file
    output_path = "/tmp/generated_strategy_pack.pptx"
    prs.save(output_path)
    return output_path
